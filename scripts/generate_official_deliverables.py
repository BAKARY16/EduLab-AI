"""Generate the editable official EduLab AI delivery pack (DOCX and PPTX)."""
from __future__ import annotations

from pathlib import Path
from datetime import date
import json

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as PptColor
from pptx.util import Inches as PInches, Pt as PPt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "deliverables" / "official"
OUT.mkdir(parents=True, exist_ok=True)
TODAY = date(2026, 7, 28).strftime("%d/%m/%Y")

GREEN = "175C44"; DARK = "14221C"; CREAM = "F7F5EF"; ORANGE = "D97836"; PALE = "E7F0E7"; GREY = "66736D"


def doc_setup(title: str, subtitle: str) -> Document:
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(.7); section.bottom_margin = Inches(.7)
    section.left_margin = Inches(.8); section.right_margin = Inches(.8)
    styles = doc.styles
    styles["Normal"].font.name = "Aptos"; styles["Normal"].font.size = Pt(10.5); styles["Normal"].font.color.rgb = RGBColor.from_string(DARK)
    for name, size, color in [("Title", 30, GREEN), ("Heading 1", 20, GREEN), ("Heading 2", 14, ORANGE), ("Heading 3", 11, GREEN)]:
        styles[name].font.name = "Aptos Display"; styles[name].font.size = Pt(size); styles[name].font.color.rgb = RGBColor.from_string(color)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.space_after = Pt(18)
    r = p.add_run("EduLab AI"); r.bold = True; r.font.size = Pt(34); r.font.color.rgb = RGBColor.from_string(GREEN)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(title); r.bold = True; r.font.size = Pt(24); r.font.color.rgb = RGBColor.from_string(DARK)
    p = doc.add_paragraph(subtitle); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].font.size = Pt(13); p.runs[0].font.color.rgb = RGBColor.from_string(GREY)
    image = ROOT / "public/images/edulab-dashboard-laptop.png"
    if image.exists():
        p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.add_run().add_picture(str(image), width=Inches(6.6))
    p = doc.add_paragraph(f"Version 1.0 — {TODAY}\nÉquipe EduLab AI — Côte d’Ivoire"); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].font.color.rgb = RGBColor.from_string(GREY)
    doc.add_page_break()
    return doc


def add_heading(doc: Document, text: str, level: int = 1):
    doc.add_heading(text, level=level)


def add_bullets(doc: Document, items):
    for item in items:
        doc.add_paragraph(item, style="List Bullet")


def add_table(doc: Document, headers, rows, widths=None):
    table = doc.add_table(rows=1, cols=len(headers)); table.alignment = WD_TABLE_ALIGNMENT.CENTER; table.style = "Light Shading Accent 1"
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]; cell.text = header; cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        for run in cell.paragraphs[0].runs: run.bold = True
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row): cells[i].text = str(value)
    return table


def add_footer(doc: Document, label: str):
    for section in doc.sections:
        p = section.footer.paragraphs[0]; p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run(f"EduLab AI — {label} — Version 1.0 — {TODAY}").font.size = Pt(8)


def project_dossier():
    doc = doc_setup("Dossier officiel du projet", "Problème, solution, architecture, gouvernance, impact et feuille de route")
    add_heading(doc, "1. Résumé exécutif")
    doc.add_paragraph("EduLab AI est une plateforme éducative ivoirienne qui regroupe cours structurés, accompagnement pédagogique, exercices, préparation au BEPC et au BAC, laboratoires virtuels et suivi des acquis. Son objectif est de rendre l’accompagnement personnalisé plus accessible, sans présenter l’intelligence artificielle comme un substitut à l’enseignant.")
    add_heading(doc, "2. Problème identifié")
    add_bullets(doc, [
        "Accès inégal à un accompagnement individualisé en dehors de la classe.",
        "Difficulté à identifier rapidement les incompréhensions et les notions à revoir.",
        "Ressources pédagogiques dispersées et qualité variable des contenus disponibles en ligne.",
        "Manque d’environnements simples pour expérimenter, s’entraîner et suivre sa progression.",
    ])
    add_heading(doc, "3. Solution proposée")
    add_table(doc, ["Composant", "Rôle", "Valeur"], [
        ["Professeur EduLab", "Questions, explications et démonstrations", "Accompagnement contextualisé"],
        ["Cours", "Sommaires et séquences pédagogiques", "Progression structurée"],
        ["Aide aux devoirs", "Analyse, tentative, indices et correction", "Apprendre sans faire à la place de l’élève"],
        ["Examens", "Entraînement et chronométrage", "Préparation BEPC/BAC"],
        ["Laboratoire", "Simulations et protocole scientifique", "Observer et expérimenter"],
        ["Résultats", "Suivi des activités et recommandations", "Repérer les priorités"],
    ])
    add_heading(doc, "4. Public cible et cas d’usage")
    add_bullets(doc, ["Élèves du collège et du lycée.", "Candidats au BEPC et au BAC.", "Enseignants souhaitant suivre les difficultés récurrentes.", "Établissements recherchant un complément numérique aux cours en présentiel."])
    add_heading(doc, "5. Parcours utilisateur")
    doc.add_paragraph("Inscription → personnalisation du niveau → tableau de bord → cours ou question → exercice/expérience → résultat enregistré → recommandation pédagogique.")
    add_heading(doc, "6. Architecture fonctionnelle et technique")
    add_table(doc, ["Couche", "Technologies", "Responsabilité"], [
        ["Interface", "Next.js 16, React 19, TypeScript, Tailwind CSS", "Expérience web responsive"],
        ["API", "FastAPI et routes Next.js", "Orchestration métier et IA"],
        ["Données", "Supabase, PostgreSQL", "Comptes, progression, ressources"],
        ["Recherche", "TF-IDF, MiniLM/FAISS, fusion de rangs", "Retrouver les passages pertinents"],
        ["Raisonnement", "OpenAI; repli extractif/local", "Produire une réponse structurée"],
        ["Recherche web", "SerpAPI et Gemini optionnels", "Complément lorsque le corpus est insuffisant"],
        ["Voix", "STT/TTS avec solutions de secours", "Questions et réponses vocales"],
    ])
    add_heading(doc, "7. Principes IA responsables")
    add_bullets(doc, [
        "Priorité au contexte du cours et aux ressources validées.",
        "Séparation entre recherche de sources et génération de la réponse finale.",
        "Signalement des limites lorsqu’une preuve suffisante n’est pas disponible.",
        "Le modèle Qwen-LoRA local reste une preuve expérimentale et non une autorité scientifique.",
        "Validation humaine obligatoire avant généralisation à des apprenants réels.",
    ])
    add_heading(doc, "8. Données")
    doc.add_paragraph("Le corpus professeur actuel comporte 252 exemples d’entraînement, 63 de validation et 54 de test. Il comprend des tâches d’explication, de plan de cours, d’exercice, de correction, de résumé et de quiz. Une partie des exemples est synthétique et doit rester clairement identifiée.")
    add_heading(doc, "9. Résultats vérifiés")
    add_table(doc, ["Évaluation", "Résultat", "Portée"], [
        ["RAG Recall@3", "83,3 %", "6 cas de référence seulement"],
        ["RAG MRR", "0,833", "6 cas de référence seulement"],
        ["Respect des filtres", "100 %", "Jeu de test RAG interne"],
        ["Tests applicatifs", "13 réussis, 1 ignoré", "Dernière exécution locale"],
        ["Génération scientifique", "En attente", "Annotations humaines nécessaires"],
    ])
    add_heading(doc, "10. Modèle économique")
    add_bullets(doc, [
        "Accès gratuit limité pour découvrir les cours et outils essentiels.",
        "Abonnement individuel pour suivi avancé, préparation intensive et usage étendu du professeur.",
        "Licences établissements avec gestion de groupes et tableaux de suivi.",
        "Partenariats de contenu et déploiements institutionnels, sous contrôle contractuel et pédagogique.",
    ])
    add_heading(doc, "11. Risques et mesures")
    add_table(doc, ["Risque", "Mesure"], [
        ["Réponse incorrecte", "RAG, verrouillage du contexte, sources, refus si insuffisant"],
        ["Contenu non autorisé", "Inventaire, statut de validation et contrôle des licences"],
        ["Données de mineurs", "Minimisation, contrôle d’accès et absence de secrets dans le dépôt"],
        ["Dépendance aux API", "Replis locaux/extractifs et supervision des fournisseurs"],
        ["Coût", "Cache, limitation des appels et routage par besoin"],
    ])
    add_heading(doc, "12. Feuille de route")
    add_bullets(doc, [
        "Court terme : valider humainement les contenus prioritaires et élargir les tests RAG.",
        "Moyen terme : pilote encadré avec enseignants et élèves consentants.",
        "Long terme : couverture nationale progressive, mode faible connexion et partenariats établissements.",
    ])
    add_heading(doc, "13. Conclusion")
    doc.add_paragraph("EduLab AI démontre la faisabilité d’un environnement éducatif intégré et contextualisé. La prochaine étape n’est pas d’ajouter davantage de promesses, mais d’augmenter la validation pédagogique, la couverture documentaire et les tests avec des utilisateurs encadrés.")
    add_footer(doc, "Dossier officiel")
    doc.save(OUT / "01_DOSSIER_OFFICIEL_EDULAB_AI.docx")


def user_deployment_guide():
    doc = doc_setup("Guide utilisateur et déploiement", "Utilisation de la plateforme, installation locale et procédures d’exploitation")
    add_heading(doc, "1. Guide rapide de l’apprenant")
    add_bullets(doc, [
        "Créer un compte, confirmer son adresse si la confirmation Supabase est activée, puis se connecter.",
        "Renseigner son niveau, sa classe, ses objectifs et ses matières prioritaires.",
        "Utiliser le tableau de bord pour reprendre une activité ou accéder aux espaces.",
        "Dans Professeur IA, sélectionner le niveau et la matière avant de poser une question précise.",
        "Dans Aide aux devoirs, fournir l’énoncé puis sa tentative avant de demander la correction.",
        "Dans Laboratoire, formuler une hypothèse, modifier une variable, relever les mesures et conclure.",
        "Dans Examens, choisir entraînement ou chronométrage, puis consulter les notions à revoir.",
    ])
    add_heading(doc, "2. Règles d’usage")
    add_bullets(doc, ["Vérifier les réponses importantes avec le cours ou l’enseignant.", "Ne jamais partager de mot de passe ou de données sensibles dans une question.", "Signaler une réponse hors sujet ou scientifiquement douteuse.", "Considérer les simulations comme des compléments, pas comme des expériences réelles."])
    add_heading(doc, "3. Prérequis techniques")
    add_table(doc, ["Élément", "Version/usage"], [["Node.js", "20 ou supérieur"], ["Python", "3.11 ou supérieur"], ["PostgreSQL/Supabase", "Base et authentification"], ["Navigateur", "Chrome, Edge ou Firefox récent"], ["GPU", "Optionnel; recommandé pour Qwen 1.5B"]])
    add_heading(doc, "4. Installation locale")
    for command in ["npm install", "Copy-Item .env.example .env.local", "npm run dev", ".venv\\Scripts\\python.exe -m uvicorn app.main:app --app-dir apps/api --port 8000", "scripts\\run_teacher_model.cmd"]:
        p = doc.add_paragraph(); r = p.add_run(command); r.font.name = "Consolas"; r.font.size = Pt(9); r.font.color.rgb = RGBColor.from_string(GREEN)
    add_heading(doc, "5. Configuration")
    doc.add_paragraph("Les variables sensibles sont renseignées dans `.env.local` et ne doivent jamais être commitées. Les clés publiées ou partagées dans une conversation doivent être révoquées.")
    add_bullets(doc, ["Supabase : URL publique, clé publiable, clé serveur et JWKS.", "API : DATABASE_URL, URL frontend et origines CORS.", "IA : fournisseur, modèle, clés OpenAI/Gemini/SerpAPI.", "Voix : fournisseur STT/TTS et identifiant de voix.", "Local : TEACHER_BASE_MODEL et TEACHER_ADAPTER_PATH."])
    add_heading(doc, "6. Vérification avant démonstration")
    add_bullets(doc, ["Ports 3000, 8000 et 8010 actifs.", "Connexion et déconnexion testées.", "Une question Professeur IA testée.", "Une expérience locale testée; connexion disponible pour PhET.", "Un examen terminé et résultat enregistré.", "Aucun secret visible à l’écran ou dans le terminal."])
    add_heading(doc, "7. Dépannage")
    add_table(doc, ["Symptôme", "Contrôle"], [["Erreur de connexion", "Variables Supabase, URL de redirection, confirmation e-mail"], ["API indisponible", "GET /api/v1/health sur le port 8000"], ["Modèle local absent", "GET /health sur le port 8010"], ["Réponse lente", "Modèle local sur CPU; utiliser GPU ou moteur distant"], ["404 laboratoire", "Vérifier la clé de simulation et la connexion PhET"], ["Build DATABASE_URL", "Charger `.env.local` avant `npm run build`"]])
    add_footer(doc, "Guide utilisateur et déploiement")
    doc.save(OUT / "02_GUIDE_UTILISATEUR_ET_DEPLOIEMENT.docx")


def validation_report():
    rag = json.loads((ROOT / "reports/rag_metrics.json").read_text(encoding="utf-8"))
    doc = doc_setup("Rapport Data, IA et validation", "Traçabilité des données, modèles utilisés, métriques et limites")
    add_heading(doc, "1. Inventaire des modèles et services")
    add_table(doc, ["Composant", "Rôle", "Statut"], [
        ["OpenAI", "Raisonnement pédagogique final", "Principal si clé/quota disponibles"],
        ["Gemini + Search", "Recherche web complémentaire", "Optionnel"],
        ["TF-IDF", "Recherche lexicale", "Baseline active"],
        ["MiniLM + FAISS", "Recherche sémantique", "Couche optionnelle"],
        ["Qwen2.5-0.5B LoRA", "Preuve locale", "Non validé pour production"],
        ["Qwen2.5-1.5B QLoRA", "Entraînement Colab prévu", "Notebook prêt"],
        ["STT/TTS", "Question et réponse vocales", "Selon fournisseur configuré"],
    ])
    add_heading(doc, "2. Jeu de données professeur")
    add_table(doc, ["Split", "Nombre"], [["Entraînement", 252], ["Validation", 63], ["Test", 54], ["Total", 369]])
    doc.add_paragraph("Les données couvrent Troisième, Terminale C et Terminale D, principalement en Mathématiques, Physique-Chimie et SVT. Elles servent à apprendre une structure de réponse pédagogique, pas à mémoriser l’intégralité du programme.")
    add_heading(doc, "3. Évaluation RAG")
    metrics = rag["metrics"]
    add_table(doc, ["Métrique", "Valeur"], [["Recall@3", f"{metrics['recall_at_3']*100:.1f} %"], ["MRR", f"{metrics['mrr']:.3f}"], ["Couverture moyenne des termes", f"{metrics['mean_term_coverage']*100:.1f} %"], ["Respect des filtres", f"{metrics['filter_accuracy']*100:.1f} %"], ["Cas de référence", rag["gold_cases"]]])
    doc.add_paragraph("Ces valeurs constituent un test de fumée, pas une preuve statistique de performance. Il faut élargir le jeu de référence et ajouter des annotations humaines.")
    add_heading(doc, "4. Protocole d’entraînement Colab")
    add_bullets(doc, ["QLoRA NF4 sur Qwen2.5-1.5B-Instruct.", "Masquage de la consigne dans la fonction de perte.", "Validation distincte, early stopping et sélection du meilleur checkpoint.", "Évaluation sur le test, courbes de loss et contrôle qualitatif.", "Export de l’adaptateur uniquement, avec métriques et limites."])
    add_heading(doc, "5. Critères de passage en production")
    add_bullets(doc, ["Corpus validé par matière et par niveau.", "Jeu d’au moins plusieurs centaines de questions humaines indépendantes.", "Exactitude scientifique mesurée par des enseignants.", "Tests de hors-sujet, hallucination, biais et refus.", "Latence et coût mesurés dans les conditions réelles.", "Procédure de signalement et de retrait d’un contenu."])
    add_heading(doc, "6. Conclusion de validation")
    doc.add_paragraph("Le projet dispose d’une chaîne Data/ML démontrable et reproductible. Le RAG présente des résultats encourageants sur un échantillon limité. Le modèle local est expérimental; il ne doit pas être présenté comme plus fiable que les données qui l’alimentent.")
    add_footer(doc, "Rapport Data, IA et validation")
    doc.save(OUT / "03_RAPPORT_DATA_IA_VALIDATION.docx")


def pitch_document():
    doc = doc_setup("Pitch officiel lié au PowerPoint", "Accroche, narration, preuves vérifiées et transitions diapositive par diapositive")
    add_heading(doc, "1. Pitch principal — environ 4 minutes")
    paragraphs = [
        ("Accroche — diapositives 1 et 2", "Qui, parmi nous, n’a jamais rêvé de voir un enfant de sa famille devenir médecin, ingénieur, chercheur ou inventeur ? Et lorsque cet enfant commence à avoir des difficultés en mathématiques ou en sciences, que faisons-nous ? Quand nous en avons les moyens, nous cherchons un répétiteur ou un maître à domicile. Mais cet accompagnement n’est pas accessible à toutes les familles, à tout moment et dans toutes les localités."),
        ("Le paradoxe numérique", "Dans le même temps, les téléphones et les tablettes occupent déjà une place importante dans l’univers des jeunes. Alors nous nous sommes posé une question simple : et si une partie de ce temps d’écran pouvait devenir un temps d’expérimentation, de compréhension et de progression ? Non pas avec un jeu éducatif basique, mais avec un véritable environnement pédagogique."),
        ("Le constat vérifié — diapositive 2", "Ce besoin ne repose pas seulement sur une impression. En 2026, le gouvernement ivoirien a reconnu un sérieux déficit d’enseignants dans les matières scientifiques et a ouvert un recrutement exceptionnel de 2 000 professeurs contractuels en mathématiques et sciences physiques. Après les épreuves, 1 800 ont été retenus : 1 300 en mathématiques et 500 en sciences physiques. La Banque mondiale rappelle également qu’une évaluation de 2019 ne trouvait que 17,2 % des élèves au niveau suffisant en mathématiques. Ces chiffres ne disent pas que les élèves sont incapables. Ils indiquent que l’accompagnement, les ressources et les pratiques doivent être renforcés."),
        ("La vision pédagogique", "Les sciences paraissent souvent abstraites lorsqu’elles sont présentées uniquement comme des formules à retenir. Elles deviennent plus accessibles lorsqu’on peut poser une question, observer une démonstration, manipuler une variable, faire une erreur et comprendre pourquoi elle est fausse. C’est précisément dans cet espace qu’EduLab AI intervient."),
        ("La solution — diapositives 3 et 4", "EduLab AI est une classe numérique conçue pour accompagner les élèves du collège et du lycée. La plateforme réunit des cours structurés, un professeur numérique capable de répondre aux questions, une aide aux devoirs fondée sur des indices progressifs, des laboratoires virtuels, une préparation au BEPC et au BAC et un suivi des acquis. L’objectif n’est pas de remplacer l’enseignant ou le répétiteur. L’objectif est de prolonger leur accompagnement lorsque l’élève se retrouve seul face à une incompréhension."),
        ("Le fonctionnement — diapositives 5 et 6", "Lorsqu’un élève pose une question, EduLab commence par rechercher un contexte pertinent dans les ressources disponibles. Si le corpus est insuffisant, des outils de recherche complémentaires peuvent être sollicités. Le moteur construit ensuite une explication adaptée au niveau sélectionné, l’affiche au tableau et peut la restituer par la voix. Cette architecture nous permet de séparer la recherche de preuves de la formulation pédagogique."),
        ("La preuve et l’honnêteté — diapositive 7", "Notre prototype fonctionne, mais nous ne prétendons pas qu’il est déjà parfait. Notre corpus professeur comprend actuellement 369 exemples répartis entre entraînement, validation et test. Sur un premier jeu interne de six questions, notre RAG obtient un Recall à trois de 83,3 %. Ce résultat est encourageant, mais l’échantillon est encore trop petit. La prochaine étape est donc la validation scientifique par des enseignants et un pilote encadré avec de vrais utilisateurs."),
        ("La démonstration — diapositive 8", "Dans la démonstration, vous verrez un élève se connecter, consulter son tableau de bord, interroger le professeur, utiliser un laboratoire ou un examen, puis retrouver son résultat et ses priorités de révision. Cela montre que nous ne présentons pas uniquement un chatbot, mais un parcours d’apprentissage complet."),
        ("L’impact et le modèle — diapositive 9", "Pour l’élève, EduLab veut rendre les sciences plus concrètes et l’erreur plus utile. Pour l’enseignant, la plateforme peut devenir un complément de suivi. Pour les familles, elle peut offrir un premier niveau d’accompagnement plus accessible. Le modèle envisagé combine un accès gratuit limité, des abonnements avancés et des licences pour les établissements."),
        ("Conclusion — diapositives 10 et 11", "Nous ne disons pas qu’EduLab résout aujourd’hui tous les défis de l’éducation scientifique. Nous montrons qu’il est possible de transformer un écran déjà présent dans la vie des jeunes en un espace où l’on questionne, expérimente et progresse. Notre ambition est simple : aider davantage d’élèves ivoiriens à ne plus subir les sciences, mais à les comprendre, les pratiquer et peut-être, demain, à en faire leur métier."),
    ]
    for title, text in paragraphs:
        add_heading(doc, title, 2); doc.add_paragraph(text)
    add_heading(doc, "2. Questions d’ouverture possibles")
    add_bullets(doc, [
        "Qui n’a jamais rêvé de voir son enfant devenir médecin, ingénieur, chercheur ou inventeur ?",
        "Combien de familles cherchent un répétiteur dès que les mathématiques deviennent difficiles ?",
        "Et si le téléphone qui distrait parfois l’élève pouvait aussi devenir son laboratoire et son espace d’entraînement ?",
        "Pourquoi les matières scientifiques restent-elles perçues comme réservées à quelques élèves, alors qu’elles peuvent s’apprendre par la pratique ?",
        "Que fait un élève à 21 heures lorsqu’il ne comprend pas une étape de son exercice et qu’aucun enseignant n’est disponible ?",
    ])
    add_heading(doc, "3. Formulations à éviter")
    add_table(doc, ["À éviter", "Formulation défendable"], [
        ["Tous les parents veulent un enfant scientifique.", "Beaucoup de familles associent les métiers scientifiques à des perspectives d’avenir; posons-le comme question d’ouverture."],
        ["Les enfants ne font que jouer sur leur téléphone.", "Les écrans sont déjà présents dans l’univers des jeunes; une partie de cet usage peut être orientée vers l’apprentissage."],
        ["Les sciences sont mal enseignées.", "Les ressources, la pratique et l’accompagnement scientifique doivent être renforcés."],
        ["L’IA remplace le manque d’enseignants.", "EduLab complète l’accompagnement et ne remplace pas l’enseignant."],
        ["Notre RAG est fiable à 83 %.", "Le Recall@3 est de 83,3 % sur seulement six cas internes; une évaluation plus large est nécessaire."],
    ])
    add_heading(doc, "4. Réponses courtes aux objections du jury")
    add_table(doc, ["Question", "Réponse recommandée"], [
        ["Pourquoi une IA alors que Mon École à la Maison existe ?", "EduLab ne remplace pas cette ressource. Il ajoute interaction, accompagnement, simulations et suivi individualisé autour de contenus traçables."],
        ["Comment éviter les fausses réponses ?", "Verrouillage du contexte, recherche de sources, séparation recherche/génération, signalement des limites et validation humaine."],
        ["Et sans Internet ?", "Le prototype dépend encore partiellement du réseau. La feuille de route prévoit cache, ressources locales et mode faible connexion."],
        ["Pourquoi entraîner Qwen ?", "Pour expérimenter un repli open source et apprendre le format pédagogique; le modèle local actuel n’est pas présenté comme l’autorité finale."],
        ["Quel est votre avantage principal ?", "L’intégration dans un même parcours : comprendre, pratiquer, expérimenter, être évalué et suivre sa progression."],
    ])
    add_heading(doc, "5. Sources vérifiées utilisées dans le pitch")
    sources = [
        "Gouvernement de Côte d’Ivoire, 22 mars 2026 — recrutement exceptionnel de 2 000 professeurs de mathématiques et sciences physiques : https://gouv.ci/actualite/recrutement-exceptionnel-de-professeurs-contractuels-de-colleges-et-lycees-de-lenseignement-secondaire-general-en-mathematiques-et-physiques-7-668-candidats-composent-pour-2-000-places-7110",
        "Gouvernement de Côte d’Ivoire, 30 mars 2026 — formation de 1 800 enseignants retenus, dont 1 300 en mathématiques et 500 en sciences physiques : https://gouv.ci/actualite/education-nationale-la-premiere-session-de-formation-des-1-800-enseignants-contractuels-de-lycees-et-colleges-recrutes-en-mathematiques-et-sciences-physiques-lancee-2093",
        "Banque mondiale, 3 février 2026 — apprentissage adaptatif et résultats d’apprentissage en Côte d’Ivoire : https://blogs.worldbank.org/en/education/adaptive-learning--a-response-to-cote-d-ivoire-s-education-chall",
        "UNICEF Côte d’Ivoire, 3 février 2023 — Mon École à la Maison et attractivité du numérique : https://www.unicef.org/cotedivoire/recits/%C2%AB-mon-ecole-%C3%A0-la-maison-%C2%BB-pour-une-%C3%A9ducation-digitale-de-qualit%C3%A9-en-c%C3%B4te-divoire",
        "UNICEF Côte d’Ivoire — Guide d’utilisation du web pour les ados et jeunes : https://www.unicef.org/cotedivoire/documents/guide-dutilisation-du-web-pour-les-ados-et-jeunes",
    ]
    for source in sources: doc.add_paragraph(source, style="List Bullet")
    add_footer(doc, "Pitch officiel")
    doc.save(OUT / "07_PITCH_OFFICIEL_LIE_AU_POWERPOINT.docx")


def ppt_textbox(slide, x, y, w, h, text, size=20, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h)); frame = box.text_frame; frame.clear(); frame.word_wrap = True; frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]; p.alignment = align; r = p.add_run(); r.text = text; r.font.name = "Aptos"; r.font.size = PPt(size); r.font.bold = bold; r.font.color.rgb = PptColor.from_string(color)
    return box


def ppt_bg(slide, color=CREAM):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = PptColor.from_string(color)


def ppt_header(slide, number, title, eyebrow="EDULAB AI"):
    ppt_textbox(slide, .65, .28, 3.5, .35, eyebrow, 9, GREEN, True)
    ppt_textbox(slide, .65, .7, 11.8, .65, title, 28, DARK, True)
    ppt_textbox(slide, 12.2, .28, .45, .3, f"{number:02d}", 9, GREY, True, PP_ALIGN.RIGHT)


def ppt_bullets(slide, items, x=.8, y=1.65, w=5.5, h=4.9, size=19):
    box = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h)); frame = box.text_frame; frame.clear(); frame.word_wrap = True
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph(); p.text = item; p.level = 0; p.font.name = "Aptos"; p.font.size = PPt(size); p.font.color.rgb = PptColor.from_string(DARK); p.space_after = PPt(13); p.text = "•  " + item


def ppt_card(slide, x, y, w, h, title, body, accent=GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(x), PInches(y), PInches(w), PInches(h)); shape.fill.solid(); shape.fill.fore_color.rgb = PptColor.from_string("FFFFFF"); shape.line.color.rgb = PptColor.from_string(PALE)
    ppt_textbox(slide, x+.25, y+.2, w-.5, .35, title, 14, accent, True)
    ppt_textbox(slide, x+.25, y+.65, w-.5, h-.82, body, 11, GREY)


def presentation():
    prs = Presentation(); prs.slide_width = PInches(13.333); prs.slide_height = PInches(7.5); blank = prs.slide_layouts[6]
    # 1
    s = prs.slides.add_slide(blank); ppt_bg(s, DARK)
    ppt_textbox(s, .75, .6, 4, .35, "PLATEFORME ÉDUCATIVE IVOIRIENNE", 10, "9DC6A9", True)
    ppt_textbox(s, .75, 1.2, 6.6, 1.3, "EduLab AI", 42, "FFFFFF", True)
    ppt_textbox(s, .75, 2.45, 5.9, 1.15, "Un accompagnement pédagogique personnalisé, structuré et accessible.", 24, "E8EFE9")
    img = ROOT / "public/images/edulab-dashboard-laptop.png"
    if img.exists(): s.shapes.add_picture(str(img), PInches(7.0), PInches(.75), width=PInches(5.7))
    ppt_textbox(s, .75, 6.65, 5, .3, "Présentation de la solution — 2026", 10, "9DC6A9")
    # 2
    s = prs.slides.add_slide(blank); ppt_bg(s); ppt_header(s, 2, "Le défi : accompagner chaque élève")
    ppt_bullets(s, ["L’accompagnement individualisé reste difficile à obtenir hors de la classe.", "Les ressources disponibles sont dispersées et de qualité inégale.", "Les difficultés sont souvent détectées trop tard.", "Les élèves ont besoin de pratiquer, expérimenter et comprendre leurs erreurs."], .8, 1.55, 6.1, 4.8, 20)
    ppt_card(s, 7.25, 1.55, 2.4, 2.0, "17,2 %", "des élèves au niveau suffisant en mathématiques dans l’évaluation 2019 citée par la Banque mondiale", ORANGE)
    ppt_card(s, 9.9, 1.55, 2.4, 2.0, "1 800", "enseignants scientifiques retenus après le recrutement exceptionnel de 2026", GREEN)
    ppt_textbox(s,7.25,4.05,5.05,1.0,"Comment renforcer l’accompagnement sans remplacer l’enseignant ?",18,DARK,True,PP_ALIGN.CENTER)
    # 3
    s = prs.slides.add_slide(blank); ppt_bg(s); ppt_header(s, 3, "La solution : une classe numérique intégrée")
    cards=[("Cours structurés","Sommaires, séquences et vérifications."),("Professeur EduLab","Questions, explications et tableau."),("Entraînement","Devoirs, exercices et examens."),("Laboratoire","Simulations et démarche scientifique."),("Suivi","Activités, maîtrise et priorités."),("Voix","Questions et réponses vocales.")]
    for i,(t,b) in enumerate(cards): ppt_card(s,.75+(i%3)*4.15,1.55+(i//3)*2.35,3.75,1.85,t,b, GREEN if i%2==0 else ORANGE)
    # 4
    s = prs.slides.add_slide(blank); ppt_bg(s); ppt_header(s, 4, "Une expérience centrée sur l’apprentissage")
    flow=[("1","Se connecter et choisir son niveau"),("2","Apprendre ou poser une question"),("3","S’entraîner et expérimenter"),("4","Recevoir une recommandation")]
    for i,(n,t) in enumerate(flow):
        x=.75+i*3.15; shape=s.shapes.add_shape(MSO_SHAPE.OVAL,PInches(x),PInches(2.0),PInches(.75),PInches(.75)); shape.fill.solid(); shape.fill.fore_color.rgb=PptColor.from_string(GREEN); shape.line.fill.background(); ppt_textbox(s,x,2.0,.75,.75,n,18,"FFFFFF",True,PP_ALIGN.CENTER); ppt_textbox(s,x-.15,3.0,2.7,1.1,t,15,DARK,True)
        if i<3: ppt_textbox(s,x+2.45,2.1,.5,.5,"→",22,ORANGE,True,PP_ALIGN.CENTER)
    # 5
    s = prs.slides.add_slide(blank); ppt_bg(s); ppt_header(s, 5, "Architecture IA : rechercher avant de répondre")
    stages=[("Question","Texte ou voix"),("Contexte","Cours actif"),("Recherche","RAG local + Web"),("Raisonnement","Moteur final"),("Restitution","Tableau + voix")]
    for i,(t,b) in enumerate(stages):
        ppt_card(s,.55+i*2.55,2.0,2.15,1.55,t,b,GREEN if i<3 else ORANGE)
        if i<4:ppt_textbox(s,2.7+i*2.55,2.45,.35,.35,"›",26,ORANGE,True,PP_ALIGN.CENTER)
    ppt_textbox(s,.8,4.35,11.7,1.0,"Ordre de confiance : documents validés → ressources éducatives autorisées → résultats web signalés → refus si les preuves sont insuffisantes.",18,GREY,False,PP_ALIGN.CENTER)
    # 6
    s = prs.slides.add_slide(blank); ppt_bg(s); ppt_header(s, 6, "Technologies et modèles utilisés")
    tech=[("Interface","Next.js · React · TypeScript"),("Backend","FastAPI · API Next.js"),("Données","Supabase · PostgreSQL"),("Recherche","TF-IDF · MiniLM · FAISS"),("IA","OpenAI · Gemini · Qwen LoRA"),("Voix","STT · TTS · secours navigateur")]
    for i,(t,b) in enumerate(tech):ppt_card(s,.8+(i%2)*6.15,1.45+(i//2)*1.65,5.65,1.3,t,b,ORANGE if i>=4 else GREEN)
    # 7
    s = prs.slides.add_slide(blank); ppt_bg(s); ppt_header(s, 7, "Données et résultats : ce qui est réellement mesuré")
    ppt_card(s,.8,1.55,3.7,2.0,"369 exemples","252 entraînement\n63 validation\n54 test",GREEN)
    ppt_card(s,4.8,1.55,3.7,2.0,"RAG Recall@3","83,3 %\nTest interne sur 6 cas",ORANGE)
    ppt_card(s,8.8,1.55,3.7,2.0,"RAG MRR","0,833\nFiltres respectés : 100 %",GREEN)
    ppt_textbox(s,.9,4.15,11.5,1.2,"Limite assumée : le modèle local actuel est une preuve technique. La validation scientifique humaine et un jeu de test plus large restent nécessaires.",18,DARK,True,PP_ALIGN.CENTER)
    # 8
    s = prs.slides.add_slide(blank); ppt_bg(s); ppt_header(s, 8, "Démonstration de la plateforme")
    img = ROOT / "public/images/edulab-ai-classroom.png"
    if img.exists(): s.shapes.add_picture(str(img), PInches(.75), PInches(1.45), width=PInches(6.15))
    ppt_bullets(s,["Connexion et personnalisation","Tableau de bord et navigation","Question au professeur","Laboratoire ou examen","Résultat et recommandation"],7.35,1.65,5.1,4.7,18)
    ppt_textbox(s,7.35,5.65,5.1,.45,"Vidéo de démonstration déjà disponible",12,ORANGE,True)
    # 9
    s = prs.slides.add_slide(blank); ppt_bg(s); ppt_header(s, 9, "Impact attendu et modèle économique")
    ppt_card(s,.8,1.45,3.7,2.15,"Impact apprenant","Compréhension, autonomie et entraînement régulier.",GREEN)
    ppt_card(s,4.8,1.45,3.7,2.15,"Impact enseignant","Visibilité sur les difficultés et complément de suivi.",GREEN)
    ppt_card(s,8.8,1.45,3.7,2.15,"Impact système","Ressources structurées et accès numérique progressif.",GREEN)
    ppt_textbox(s,.8,4.35,11.7,.5,"Freemium individuel  •  Abonnement avancé  •  Licences établissements  •  Partenariats institutionnels",18,ORANGE,True,PP_ALIGN.CENTER)
    # 10
    s = prs.slides.add_slide(blank); ppt_bg(s); ppt_header(s, 10, "Prochaines étapes")
    ppt_bullets(s,["Valider les contenus prioritaires avec des enseignants.","Élargir le jeu d’évaluation RAG et scientifique.","Lancer un pilote encadré avec consentement.","Optimiser le mode faible connexion et le coût IA.","Mesurer l’impact pédagogique avant généralisation."],.9,1.5,7.2,4.9,20)
    ppt_card(s,8.6,1.7,3.7,3.4,"Notre ambition","Faire d’EduLab AI un complément pédagogique crédible, accessible et adapté au contexte ivoirien.",ORANGE)
    # 11
    s = prs.slides.add_slide(blank); ppt_bg(s,DARK)
    ppt_textbox(s,1.0,1.25,11.3,.7,"Merci",36,"FFFFFF",True,PP_ALIGN.CENTER)
    ppt_textbox(s,1.2,2.15,10.9,1.1,"EduLab AI — apprendre, pratiquer et progresser avec un accompagnement contextualisé.",22,"DDE9E0",False,PP_ALIGN.CENTER)
    ppt_textbox(s,1.2,4.65,10.9,.6,"Questions / Démonstration",18,"D99A6F",True,PP_ALIGN.CENTER)
    # speaker notes are delivered separately for portability
    prs.save(OUT / "04_PRESENTATION_EDULAB_AI.pptx")


if __name__ == "__main__":
    project_dossier(); user_deployment_guide(); validation_report(); pitch_document(); presentation()
    print(f"Livrables générés dans {OUT}")
